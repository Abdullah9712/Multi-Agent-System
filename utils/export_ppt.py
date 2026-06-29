from pptx import Presentation


def create_ppt(data, output_path="output/report.pptx"):
    prs = Presentation()

    # Slide 1 - Title
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = data.get("topic", "")
    slide.placeholders[1].text = "AI Research Report"

    # Slide 2 - Summary
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Executive Summary"
    slide.placeholders[1].text = data.get("executive_summary", "")

    # Slide 3 - Overview
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Overview"
    slide.placeholders[1].text = "\n".join(data.get("overview", []))

    # Slide 4 - Key Concepts
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Key Concepts"

    concepts_text = ""
    for item in data.get("key_concepts", []):
        if isinstance(item, dict):
            concepts_text += f"{item['term']}: {item['description']}\n"
        else:
            concepts_text += item + "\n"

    slide.placeholders[1].text = concepts_text

    prs.save(output_path)

    return output_path